# ==============================================================================
# --- 1. IMPORTS
# ==============================================================================
#
import os
from dotenv import load_dotenv
import torch
import requests
import traceback
from flask import Flask, request
from collections import defaultdict
import re
from datetime import datetime, timedelta
from io import BytesIO
import tempfile
import numpy as np
import whisper 
import base64
import moviepy
import moviepy.editor
from fpdf import FPDF
from pydub import AudioSegment
from moviepy.editor import ImageClip, concatenate_videoclips, AudioFileClip
from gtts import gTTS
from PIL import Image, ImageDraw
import json
import threading



# --- AI & Machine Learning ---
import google.generativeai as genai
from sklearn.metrics.pairwise import cosine_similarity
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.document_loaders import UnstructuredFileLoader
from langchain_community.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from diffusers import DiffusionPipeline
from diffusers.utils import export_to_video

# --- Google Cloud APIs ---
import google.auth.transport.requests
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaFileUpload

# --- Utility & File Generation ---
from fpdf import FPDF
from pydub import AudioSegment
from gtts import gTTS
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches

#
# ---------------------------------
# --- MODEL LOADING (at startup)
# ---------------------------------
#

print("Loading Whisper model...")
whisper_model = whisper.load_model("base")
print("Whisper model loaded successfully.")
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"

#
# ---------------------------------
# --- ENVIRONMENT & CONFIG
# ---------------------------------
#
app = Flask(__name__)

load_dotenv()
# --- WhatsApp & Meta Config ---
VERIFY_TOKEN = os.getenv("VERIFY_TOKEN")
ACCESS_TOKEN = os.getenv("ACCESS_TOKEN")
PHONE_NUMBER_ID = os.getenv("PHONE_NUMBER_ID")
TEMPLATE_NAME = os.getenv("TEMPLATE_NAME")

# Google & GenAI
GENAI_API_KEY = os.getenv("GENAI_API_KEY")
GOOGLE_FORMS_SERVICE_ACCOUNT_KEY_PATH = os.getenv("GOOGLE_FORMS_SERVICE_ACCOUNT_KEY_PATH")
GOOGLE_FORMS_TARGET_FOLDER_ID = os.getenv("GOOGLE_FORMS_TARGET_FOLDER_ID")
CLASSROOM_COURSE_ID = os.getenv("CLASSROOM_COURSE_ID")

print("Verify Token:", VERIFY_TOKEN)
print("Configs loaded successfully")


# --- List of users to notify on startup ---
student_phone_numbers = [
    "916379613654",
    "918870420449",
]

# --- Google AI & API Config ---
genai.configure(api_key=GENAI_API_KEY)
GOOGLE_FORMS_SERVICE_ACCOUNT_KEY_PATH = 'sahayak-465916-8bf5ddce5515.json'
GOOGLE_FORMS_TARGET_FOLDER_ID = '11vL5AgJiLYbX6fgAJFoMEjea9jiR1WMk'
CLASSROOM_COURSE_ID = '791255014049'

# --- State & Memory Management ---
user_states = defaultdict(str)
user_temp_data = defaultdict(dict)
user_memory = defaultdict(list)
MAX_HISTORY = 50

# Initialize video_pipe as None
video_pipe = None
try:
    print("Loading Text-to-Video model (this may take a while)...")
    video_pipe = DiffusionPipeline.from_pretrained("damo-vilab/text-to-video-ms-1.7b", torch_dtype=torch.float16, variant="fp16")
    if torch.cuda.is_available():
        video_pipe = video_pipe.to("cuda")
        print("Text-to-Video model moved to GPU.")
    else:
        print("WARNING: No GPU detected. Video generation will be very slow on CPU.")
except Exception as e:
    print(f"CRITICAL: Failed to load Text-to-Video model. Video features will be disabled. Error: {e}")


# --- Intent Classification Examples & Thresholds ---
OS_EXAMPLES = [
    "What is an operating system?", "Explain process scheduling.", "Tell me about deadlocks in OS.",
    "How does virtual memory work?", "What are the types of file systems?", "Describe the role of the kernel.",
    "What is multitasking?", "How does memory management happen in an OS?", "What is CPU scheduling?",
    "Explain paging.", "What is a semaphore in OS?", "Discuss thread management.",
    "What are system calls?", "Explain storage management.", "What is a logical address space?",
    "What is a process control block?", "Describe context switching.", "Explain segmentation.",
    "What is I/O management?", "How does buffering work in OS?", "What is thrashing?",
    "Explain mutual exclusion.", "What are producers-consumers problem?", "Describe inter-process communication.",
    "What is fragmentation?", "What is the OS syllabus?", "Can you provide details on the operating system syllabus?",
    "Describe my OS syllabus.", "Give me an overview of the OS course content.",
    "What topics are covered in the OS subject?", "Show me the syllabus for Operating Systems.",
    "List the units in the OS syllabus?", "What are the main topics in the OS syllabus?",
    "Can you summarize the OS syllabus for me?", "What is the structure of the OS syllabus?",
    "What are the key areas in the OS syllabus?", "How is the OS syllabus organized?",
    "Tell me about the course outline for OS."
]
FORM_GENERATION_EXAMPLES = [
    "Generate a G form link for 10 MCQ questions in CPU scheduling.", "Create a quiz on memory management.",
    "Make 5 multiple choice questions about deadlocks.", "I need a short quiz on OS concepts.",
    "Can you generate a Google Form for me?", "Form for process synchronization with 7 questions.",
    "Build a form for 20 questions on operating systems.", "Generate a quiz form about processes."
]
CLASSROOM_ASSIGNMENT_EXAMPLES = [
    "Create an assignment 'CPU Scheduling Quiz' for 'Operating Systems Class'.",
    "Post a new assignment in OS class about memory management.", "Make an assignment for the operating systems course.",
    "Create a quiz in Classroom for Deadlocks.", "New assignment for my OS students on virtual memory.",
    "Assign a new quiz on processes.", "Create homework on threading for OS class."
]
WORKSHEET_GENERATION_EXAMPLES = [
    "Generate a worksheet on deadlocks.", "Create a PDF worksheet for memory management.",
    "Give me a short answer worksheet on CPU scheduling.", "Make a worksheet with 5 fill-in-the-blanks on virtual memory.",
    "Can you generate a practice sheet for OS concepts?", "Generate a PDF for operating systems review.",
    "Generate 10 multiple choice questions worksheet on processes.", "Create an MCQ worksheet for file systems.",
    "Give me 5 short answer questions on concurrency.", "Generate a short answer worksheet about synchronization.",
    "Make 3 long answer questions on memory allocation strategies.", "Give me an essay question worksheet on virtual memory benefits.",
    "Generate 2 numerical questions on CPU scheduling algorithms.", "Create a worksheet with calculation problems about page replacement."
]
def send_whatsapp_document(to, file_bytes, filename):
    media_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/media"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    files = {'file': (filename, file_bytes, 'application/pdf'), 'messaging_product': (None, 'whatsapp')}
    try:
        media_res = requests.post(media_url, headers=headers, files=files)
        media_res.raise_for_status()
        media_id = media_res.json().get("id")
        if not media_id: return "failure"
        
        msg_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
        msg_payload = {"messaging_product": "whatsapp", "to": to, "type": "document", "document": {"id": media_id, "filename": filename}}
        msg_res = requests.post(msg_url, headers=headers, json=msg_payload)
        msg_res.raise_for_status()
        return "success"
    except Exception as e:
        print(f"Failed to send document: {e}")
        return "failure"

embeddings_model_for_classification = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
os_example_embeddings = None
form_generation_example_embeddings = None
classroom_assignment_example_embeddings = None
worksheet_generation_example_embeddings = None

SIMILARITY_THRESHOLD_OS = 0.65
SIMILARITY_THRESHOLD_FORM = 0.70
SIMILARITY_THRESHOLD_CLASSROOM = 0.70
SIMILARITY_THRESHOLD_WORKSHEET = 0.75

# --- API Service Clients (Initialized later) ---
forms_service = None
drive_service = None
classroom_service = None
speech_client = None
creds = None
calendar_service = None

#
# ---------------------------------
# --- API INITIALIZATION
# ---------------------------------
#
SCOPES = [
    'https://www.googleapis.com/auth/classroom.courses',
    'https://www.googleapis.com/auth/classroom.announcements',
    'https://www.googleapis.com/auth/drive.file',
    'https://www.googleapis.com/auth/forms.body',
    'https://www.googleapis.com/auth/calendar.events'
]

def init_google_apis():
    """Authenticates user and initializes Google API clients."""
    global classroom_service, drive_service, forms_service, calendar_service, creds
    if os.path.exists('token.json'):
        creds = Credentials.from_authorized_user_file('token.json', SCOPES)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(google.auth.transport.requests.Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file('credentials.json', SCOPES)
            creds = flow.run_local_server(port=0)
        with open('token.json', 'w') as token:
            token.write(creds.to_json())
    
    try:
        drive_service = build('drive', 'v3', credentials=creds)
        classroom_service = build('classroom', 'v1', credentials=creds)
        forms_service = build('forms', 'v1', credentials=creds)
        calendar_service = build('calendar', 'v3', credentials=creds) # <-- Initialize Calendar
        print("Google APIs initialized successfully.")
    except Exception as e:
        print(f"Error initializing Google APIs: {e}")


# Add this function to your WHATSAPP & MEMORY FUNCTIONS section
def send_start_template(to_number):
    """Sends the approved 'start' template message to a phone number."""
    url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}", "Content-Type": "application/json"}
    payload = {
        "messaging_product": "whatsapp",
        "to": to_number,
        "type": "template",
        "template": {
            "name": "start_conversation_prompt", # Or your approved template name
            "language": {
                "code": "en"
            }
        }
    }
    try:
        response = requests.post(url, headers=headers, json=payload)
        response.raise_for_status()
        print(f"Successfully sent 'START' template to {to_number}.")
    except requests.exceptions.RequestException as e:
        print(f"Failed to send template to {to_number}: {e}")
        print(f"Response: {response.text}")

# ---------------------------------
# --- WHATSAPP & MEMORY FUNCTIONS
# ---------------------------------
#
def append_to_memory(user_id, role, content):
    """
    Appends a message to a user's conversation history in the correct
    format for the Gemini API.
    """
    # The Gemini API expects "user" and "model" as roles.
    # This maps your internal "assistant" role to "model".
    api_role = "model" if role == "assistant" else "user"
    
    # This uses the 'parts' key required by the API instead of 'content'.
    user_memory[user_id].append({"role": api_role, "parts": [str(content)]})
    
    if len(user_memory[user_id]) > MAX_HISTORY:
        user_memory[user_id].pop(0)

def send_whatsapp_message(to, message):
    url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}", "Content-Type": "application/json"}
    payload = {"messaging_product": "whatsapp", "to": to, "text": {"body": message}}
    try:
        res = requests.post(url, headers=headers, json=payload)
        res.raise_for_status()
        print(f"WhatsApp text message sent to {to}: {message[:75]}...")
    except Exception as e:
        print(f"Failed to send WhatsApp text message: {e}")

def send_whatsapp_document(to, file_bytes, filename):
    media_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/media"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    files_and_data = {
        'file': (filename, file_bytes, 'application/pdf'),
        'messaging_product': (None, 'whatsapp'), 'type': (None, 'document')
    }
    try:
        media_res = requests.post(media_url, headers=headers, files=files_and_data)
        media_res.raise_for_status()
        media_id = media_res.json().get("id")
        if not media_id:
            raise Exception("Failed to get media ID from upload response.")
    except Exception as e:
        print(f"Failed to upload PDF to WhatsApp: {e}")
        return
    
    message_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
    message_headers = {"Authorization": f"Bearer {ACCESS_TOKEN}", "Content-Type": "application/json"}
    message_payload = {"messaging_product": "whatsapp", "to": to, "type": "document", "document": {"id": media_id, "filename": filename}}
    try:
        msg_res = requests.post(message_url, headers=message_headers, json=message_payload)
        msg_res.raise_for_status()
        print(f"WhatsApp document sent to {to}: {filename}")
    except Exception as e:
        print(f"Failed to send WhatsApp document message: {e}")

def send_menu_message(to, text, options):
    """Sends an interactive List Message."""
    url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}", "Content-Type": "application/json"}
    
    rows = [{"id": f"option_{i}", "title": option} for i, option in enumerate(options)]
    
    payload = {
        "messaging_product": "whatsapp",
        "to": to,
        "type": "interactive",
        "interactive": {
            "type": "list",
            "header": {"type": "text", "text": "Main Menu"},
            "body": {"text": text},
            "footer": {"text": "Please choose an option"},
            "action": {
                "button": "Choose an Option",
                "sections": [{"title": "Available Actions", "rows": rows}]
            }
        }
    }
    try:
        # ✅ CORRECTED: Explicitly dump the payload to a JSON string
        requests.post(url, headers=headers, data=json.dumps(payload)).raise_for_status()
        print(f"List Menu sent to {to}.")
    except Exception as e:
        print(f"Failed to send List Menu: {e}")

def send_whatsapp_video(to, video_bytes, filename):
    media_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/media"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    files = {'file': (filename, video_bytes, 'video/mp4'), 'messaging_product': (None, 'whatsapp')}
    try:
        media_res = requests.post(media_url, headers=headers, files=files)
        media_res.raise_for_status()
        media_id = media_res.json().get("id")
        if not media_id: return "failure"
        
        msg_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
        msg_payload = {"messaging_product": "whatsapp", "to": to, "type": "video", "video": {"id": media_id}}
        msg_res = requests.post(msg_url, headers=headers, json=msg_payload)
        msg_res.raise_for_status()
        return "success"
    except Exception as e:
        print(f"Failed to send video: {e}")
        return "failure"
    
def handle_audio_task(sender, media_id):
    """Worker function to handle audio processing in a background thread."""
    print(f"Starting audio task for {sender}")
    send_whatsapp_message(sender, "Got your voice message, listening... 🎤")
    if media_url := get_media_url(media_id):
        if audio_bytes := download_media_file(media_url):
            if transcribed_text := transcribe_audio(audio_bytes):

                process_message(sender, transcribed_text)
            else:
                send_whatsapp_message(sender, "Sorry, I couldn't understand your audio.")
        else:
            send_whatsapp_message(sender, "Sorry, I couldn't download your voice message.")
    print(f"Audio task for {sender} finished.")
    
def handle_image_task(sender, media_id, prompt):
    """Worker function to handle image processing in a background thread."""
    print(f"Starting image task for {sender}")
    send_whatsapp_message(sender, "Analyzing your image... 🖼️")
    if media_url := get_media_url(media_id):
        if image_bytes := download_media_file(media_url):
            if 'podcast' in prompt.lower() or 'audio' in prompt.lower():
                ocr_prompt = "Extract all the text from this image. Do not summarize or explain, just provide the raw text."
                ocr_response = query_gemini_vision(ocr_prompt, image_bytes)
                
                if ocr_response.get("source") == "gemini_vision":
                    extracted_text = ocr_response["result"]
                    if audio_bytes := generate_voiceover(extracted_text):
                        send_whatsapp_audio(sender, audio_bytes, "podcast.mp3")
                    else:
                        send_whatsapp_message(sender, "I extracted the text, but couldn't create the audio.")
                else:
                    send_whatsapp_message(sender, ocr_response["result"])
            else:
                response = query_gemini_vision(prompt, image_bytes)
                send_whatsapp_message(sender, response["result"])
        else:
            send_whatsapp_message(sender, "Sorry, I couldn't download the image.")
    print(f"Image task for {sender} finished.")




# ---------------------------------
# --- VOICE INPUT HELPER FUNCTIONS
# ---------------------------------
#
def get_media_url(media_id):
    """Retrieves the downloadable URL for a piece of media from WhatsApp."""
    url = f"https://graph.facebook.com/v19.0/{media_id}/"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    try:
        res = requests.get(url, headers=headers)
        res.raise_for_status()
        return res.json().get("url")
    except Exception as e:
        print(f"Failed to get media URL for ID {media_id}: {e}")
        return None

def download_media_file(media_url):
    """Downloads the media file from the given URL."""
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    try:
        res = requests.get(media_url, headers=headers)
        res.raise_for_status()
        return res.content
    except Exception as e:
        print(f"Failed to download media from URL: {e}")
        return None

def transcribe_audio(audio_bytes):
    """Transcribes audio bytes using the self-hosted Whisper model."""
    print("Transcribing audio with local Whisper model...")
    
    temp_file_path = None
    try:
        # Create a temporary file, get its path, and then close it
        # so that FFmpeg can access it without a permission error.
        with tempfile.NamedTemporaryFile(suffix=".ogg", delete=False) as temp_file:
            temp_file.write(audio_bytes)
            temp_file_path = temp_file.name
        
        # Load the audio from the file path to get a NumPy array
        audio_np_array = whisper.load_audio(temp_file_path)

        # Transcribe the NumPy array
        result = whisper_model.transcribe(audio_np_array, fp16=False)
        
        transcript = result.get("text", "")
        print(f"Whisper transcription successful: '{transcript}'")
        return transcript

    except Exception as e:
        print(f"Error during local Whisper transcription: {e}")
        traceback.print_exc()
        return None
    finally:
        # Manually delete the temporary file after we're done with it
        if temp_file_path and os.path.exists(temp_file_path):
            os.remove(temp_file_path)
    
    
#
# ---------------------------------
# --- RAG & VECTOR DB FUNCTIONS
# ---------------------------------
#
def build_os_vector_index():
    print("Loading documents for OS vector index...")
    loaders = [
        UnstructuredFileLoader("data/OS_Syllabus.pdf"),
        UnstructuredFileLoader("data/OS_Unit_I_Overview.pptx"),
        UnstructuredFileLoader("data/OS_Unit_II_Process_Management.pptx"),
        UnstructuredFileLoader("data/OS_Unit_II_Process_Scheduling.pptx"),
        UnstructuredFileLoader("data/OS_Unit_III_Deadlocks.pptx"),
        UnstructuredFileLoader("data/OS_Unit_IV_Memory_management.pptx"),
        UnstructuredFileLoader("data/OS_Unit_IV_Virtual_Memory.pptx"),
        UnstructuredFileLoader("data/OS_Unit_V_File_System_Implementation.pptx"),
        UnstructuredFileLoader("data/OS_Unit_V_Storage_management.pptx"),
    ]
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    documents = []
    for loader in loaders:
        try:
            raw_docs = loader.load()
            split_docs = text_splitter.split_documents(raw_docs)
            documents.extend(split_docs)
            print(f"Loaded and split {len(raw_docs)} documents from {loader.file_path}")
        except Exception as e:
            print(f"Error loading {loader.file_path}: {e}")
    if not documents:
        print("No documents loaded, FAISS index will be empty.")
        return
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    db = FAISS.from_documents(documents, embeddings)
    db.save_local("vector_index/os_docs")
    print("OS Vector Index built successfully.")

def query_os_subject(question):
    try:
        embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        db = FAISS.load_local("vector_index/os_docs", embeddings, allow_dangerous_deserialization=True)
        retriever = db.as_retriever(search_kwargs={"k": 1})
        docs = retriever.invoke(question)
        if not docs or all(len(doc.page_content.strip()) == 0 for doc in docs):
            return {"result": "No relevant content found in OS database.", "source": "fallback_needed"}
        context = "\n\n".join([doc.page_content for doc in docs])
        sources = [os.path.basename(doc.metadata['source']) for doc in docs if 'source' in doc.metadata]
        source_info = "\n\n(Source: " + ", ".join(list(set(sources))) + ")" if sources else ""
        prompt = f"Based on the following context, answer the question accurately and concisely. If the context does not contain enough information, state that directly:\n\nContext:\n{context}\n\nQuestion: {question}\n\nAnswer:"
        model = genai.GenerativeModel('gemini-1.5-flash')
        response = model.generate_content(prompt)
        final_answer = response.text + source_info
        return {"result": final_answer, "source": "db"}
    except Exception as e:
        print("Error during OS subject query:", e)
        traceback.print_exc()
        return {"result": "Error during database lookup.", "source": "fallback_needed"}

def upload_file_to_drive(file_path):
    if not drive_service:
        print("Google Drive API not initialized.")
        return None
    try:
        file_metadata = {'name': os.path.basename(file_path)}
        media = MediaFileUpload(file_path, mimetype='application/pdf')
        file = drive_service.files().create(body=file_metadata, media_body=media, fields='id').execute()
        print(f"File uploaded to Drive with ID: {file.get('id')}")
        return file.get('id')
    except Exception as e:
        print(f"Error uploading file to Drive: {e}")
        return None
    
def post_announcement(course_id, text, drive_file_id):
    if not classroom_service:
        return {"result": "Google Classroom API not initialized.", "source": "api_error"}
    
    announcement = {
        'text': text,
        'materials': [{'driveFile': {'driveFile': {'id': drive_file_id}, 'shareMode': 'VIEW'}}],
        'state': 'PUBLISHED'
    }
    try:
        announcement = classroom_service.courses().announcements().create(courseId=course_id, body=announcement).execute()
        return {"result": f"Successfully posted announcement in Google Classroom!", "source": "classroom_success"}
    except HttpError as api_error:
        print(f"Google Classroom API Error: {api_error.content}")
        return {"result": "Couldn't post announcement. Please check permissions and Course ID.", "source": "api_error"}

    
def post_assignment(course_id, title, drive_file_id):
    if not classroom_service:
        return {"result": "Google Classroom API not initialized.", "source": "api_error"}
    coursework = {
        'title': title,
        'materials': [{'driveFile': {'driveFile': {'id': drive_file_id}, 'shareMode': 'VIEW'}}],
        'workType': 'ASSIGNMENT',
        'state': 'PUBLISHED'
    }
    try:
        assignment = classroom_service.courses().courseWork().create(courseId=course_id, body=coursework).execute()
        return {"result": f"Successfully created assignment '{assignment.get('title')}' in Google Classroom!", "source": "classroom_success"}
    except HttpError as api_error:
        print(f"Google Classroom API Error: {api_error.content}")
        return {"result": "Couldn't create Classroom assignment. Please check permissions and Course ID.", "source": "api_error"}


def query_gemini(question, history):
    try:
        model = genai.GenerativeModel('gemini-1.5-flash')
        chat = model.start_chat(history=history)
        response = chat.send_message(question)
        return {"result": response.text}
    except Exception as e:
        print("Gemini general query error:", e)
        return {"result": "Sorry, I encountered an issue while processing your general question."}
    
def generate_video_script(topic, history):
    print(f"Generating video script for topic: {topic}")
    prompt = f"Create a short, simple video script explaining the basics of {topic} for a beginner. The script should have 3 key points. For each key point, provide a narration sentence and a simple visual description on a new line starting with 'VISUAL:'."
    try:
        model = genai.GenerativeModel('gemini-1.5-flash')
        chat = model.start_chat(history=history)
        response = chat.send_message(prompt)
        return response.text
    except Exception as e:
        print(f"Error generating video script: {e}")
        return None

def generate_images_for_script(script):
    print("Generating images for script...")
    image_paths = []
    visual_prompts = re.findall(r'VISUAL:\s*(.*)', script)
    for i, prompt in enumerate(visual_prompts):
        try:
            img = Image.new('RGB', (1280, 720), color = 'darkblue')
            d = ImageDraw.Draw(img)
            try:
                font = ImageFont.truetype("arial.ttf", 40)
            except IOError:
                font = ImageFont.load_default()
            
            # Simple text wrapping
            lines = []
            words = prompt.split()
            current_line = ""
            for word in words:
                if len(current_line + " " + word) < 50:
                    current_line += " " + word
                else:
                    lines.append(current_line)
                    current_line = word
            lines.append(current_line)

            y_text = 300
            for line in lines:
                d.text((100, y_text), line.strip(), font=font, fill=(255,255,255))
                y_text += 50

            img_path = f"temp_image_{i}.png"
            img.save(img_path)
            image_paths.append(img_path)
        except Exception as e:
            print(f"Error creating image {i}: {e}")
    return image_paths

def generate_video_with_modelscope(prompt, output_path):
    """Generates a video from a text prompt using the ModelScope model."""
    print(f"Generating video for prompt: '{prompt}'...")
    try:
        video_frames = video_pipe(prompt, num_inference_steps=25, num_frames=16).frames
        export_to_video(video_frames, output_path)
        return True
    except Exception as e:
        print(f"Error generating video with ModelScope: {e}")
        return False
    
def generate_voiceover(text):
    """Generates a voiceover from text and speeds it up."""
    print("Generating voiceover...")
    try:
        # Generate the initial audio
        tts = gTTS(text=text, lang='en', slow=False)
        fp = BytesIO()
        tts.write_to_fp(fp)
        fp.seek(0)
        
        # Load the audio with pydub and speed it up
        audio = AudioSegment.from_file(fp, format="mp3")
        sped_up_audio = audio.speedup(playback_speed=1.25)
        
        # Export the sped-up audio to a new BytesIO object
        output_fp = BytesIO()
        sped_up_audio.export(output_fp, format="mp3")
        output_fp.seek(0)
        
        return output_fp.read()
    except Exception as e:
        print(f"Error generating voiceover: {e}")
        return None
    
def create_calendar_event(title, due_date):
    if not calendar_service: return
    try:
        event = {
            'summary': f'Due: {title}',
            'description': 'An assignment is due.',
            'start': {'date': due_date.isoformat()},
            'end': {'date': due_date.isoformat()},
        }
        calendar_service.events().insert(calendarId='primary', body=event).execute()
        print("Calendar event created.")
    except Exception as e:
        print(f"Error creating calendar event: {e}")

def post_assignment(course_id, title, drive_file_id, due_date=None):
    if not classroom_service: return {"result": "Classroom API not initialized."}
    try:
        coursework = {'title': title, 'workType': 'ASSIGNMENT', 'state': 'PUBLISHED'}
        if drive_file_id:
             coursework['materials'] = [{'driveFile': {'driveFile': {'id': drive_file_id}, 'shareMode': 'VIEW'}}]
        if due_date:
            coursework['dueDate'] = {'year': due_date.year, 'month': due_date.month, 'day': due_date.day}
            coursework['dueTime'] = {'hours': 23, 'minutes': 59, 'seconds': 0}
            
        classroom_service.courses().courseWork().create(courseId=course_id, body=coursework).execute()
        return {"result": f"Successfully created assignment '{title}' in Google Classroom!"}
    except Exception as e:
        return {"result": f"Error creating assignment: {e}"}

def compile_video(image_paths, voiceover_path, output_path):
    print("Compiling video...")
    try:
        clips = [ImageClip(m).set_duration(5) for m in image_paths]
        if not clips: return False
        
        video = concatenate_videoclips(clips, method="compose")
        audio = AudioFileClip(voiceover_path)
        if audio.duration > video.duration:
            audio = audio.subclip(0, video.duration)
        video = video.set_audio(audio)
        
        video.write_videofile(output_path, fps=24, codec='libx264', audio_codec='aac')
        return True
    except Exception as e:
        print(f"Error compiling video: {e}")
        return False


#
# ---------------------------------
# --- CONTENT GENERATION FUNCTIONS
# ---------------------------------
#
def generate_mcq_questions_text(topic, num_questions, history):
    print(f"Attempting to generate {num_questions} TEXT MCQs on topic: '{topic}'")
    context, source_info = "", ""
    # Simplified RAG check for brevity
    topic_embedding = embeddings_model_for_classification.embed_query(topic)
    topic_embedding_reshaped = np.array(topic_embedding).reshape(1, -1)
    similarities = cosine_similarity(topic_embedding_reshaped, os_example_embeddings)
    if np.max(similarities) > SIMILARITY_THRESHOLD_OS:
        try:
            embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
            db = FAISS.load_local("vector_index/os_docs", embeddings, allow_dangerous_deserialization=True)
            retriever = db.as_retriever(search_kwargs={"k": 3})
            docs = retriever.invoke(topic)
            if docs:
                context = "\n\n".join([doc.page_content for doc in docs])
                unique_sources = list(set([os.path.basename(d.metadata['source']) for d in docs if 'source' in d.metadata]))
                if unique_sources:
                    source_info = f"\n\n(Context from: {', '.join(unique_sources)})"
        except Exception as e:
            print(f"Error retrieving RAG context for MCQ: {e}")
    mcq_prompt = f"Generate {num_questions} ONLY multiple-choice questions on the topic of '{topic}'. For each question, provide 4 options (A, B, C, D) and clearly indicate the correct answer. Format each as: 'Question Number. Question Text\\n A) Option A\\n B) Option B\\n C) Option C\\n D) Option D\\n Correct Answer: X'. Provide ONLY the questions and answers."
    if context:
        mcq_prompt = f"Based on the following context:\n\n{context}\n\n" + mcq_prompt
    try:
        model = genai.GenerativeModel('gemini-1.5-flash')
        chat = model.start_chat(history=history)
        response = chat.send_message(mcq_prompt)
        return {"result": response.text + source_info, "source": "generated_form_questions_text"}
    except Exception as e:
        print(f"Error generating TEXT MCQs for '{topic}': {e}")
        return {"result": "Sorry, I couldn't generate the quiz questions.", "source": "error"}

def create_google_form_mcq(title, mcq_questions_text):
    """Creates a Google Form quiz from a string of MCQ questions."""
    if not forms_service:
        return {"result": "Google Forms API not initialized.", "source": "api_error"}

    try:
        form_metadata = {'info': {'title': title}}
        form = forms_service.forms().create(body=form_metadata).execute()
        form_id = form['formId']
        
        update = {"requests": [{"updateSettings": {"settings": {"quizSettings": {"isQuiz": True}}, "updateMask": "quizSettings"}}]}
        forms_service.forms().batchUpdate(formId=form_id, body=update).execute()

        question_requests = []
        # A more robust way to split into question blocks
        question_blocks = re.split(r'\n\s*\d+\.\s*', mcq_questions_text)
        
        # Filter out any empty strings or introductory text
        question_blocks = [b.strip() for b in question_blocks if b.strip() and "Instructions:" not in b and "Worksheet" not in b]

        for i, q_block in enumerate(question_blocks):
            lines = [line.strip() for line in q_block.split('\n') if line.strip()]
            if not lines:
                continue

            question_text = lines[0]
            options = []
            correct_answer_text = ""

            for line in lines[1:]:
                if re.match(r'^[a-dA-D][\.\)]', line):
                    options.append(line[2:].strip())
                elif "Correct Answer:" in line:
                    correct_answer_match = re.search(r'Correct Answer:\s*([a-dA-D])', line, re.IGNORECASE)
                    if correct_answer_match:
                        correct_answer_letter = correct_answer_match.group(1).upper()
                        correct_answer_index = ord(correct_answer_letter) - ord('A')
                        if 0 <= correct_answer_index < len(options):
                            correct_answer_text = options[correct_answer_index]

            if not question_text or len(options) < 2 or not correct_answer_text:
                print(f"Skipping malformed block: {q_block[:50]}...")
                continue

            choices = [{"value": opt} for opt in options]
            
            question_requests.append({
                "createItem": {
                    "item": {
                        "title": question_text,
                        "questionItem": {
                            "question": {
                                "required": True,
                                "choiceQuestion": {"type": "RADIO", "options": choices},
                                "grading": {"pointValue": 1, "correctAnswers": {"answers": [{"value": correct_answer_text}]}}
                            }
                        }
                    }, "location": {"index": i}
                }
            })

        if question_requests:
            forms_service.forms().batchUpdate(formId=form_id, body={"requests": question_requests}).execute()

        return {"result": f"I've created a Google Form quiz for you! Access it here: {form['responderUri']}", "source": "google_form_created"}

    except Exception as e:
        print(f"Error creating Google Form: {e}")
        traceback.print_exc()
        return {"result": "Sorry, I encountered an error while creating the Google Form.", "source": "error"}


def create_classroom_assignment(course_id, title, description_html, max_points=100):
    if not classroom_service:
        return {"result": "Google Classroom API not initialized.", "source": "api_error"}
    assignment_body = {'title': title, 'description': description_html, 'maxPoints': max_points, 'workType': 'ASSIGNMENT', 'state': 'PUBLISHED'}
    try:
        coursework = classroom_service.courses().courseWork().create(courseId=course_id, body=assignment_body).execute()
        return {"result": f"Successfully created assignment '{title}' in Google Classroom!", "source": "classroom_success"}
    except HttpError as api_error:
        print(f"Google Classroom API Error: {api_error.content}")
        return {"result": f"Couldn't create Classroom assignment due to an API error. Please check permissions and Course ID.", "source": "api_error"}
    except Exception as e:
        print(f"General error creating Classroom assignment: {e}")
        return {"result": "An unexpected error occurred while creating the Classroom assignment.", "source": "error"}




def generate_worksheet_content_text(topic, num_items, worksheet_type, history):
    """
    Generates worksheet questions AND answers, separated by a unique string.
    """
    print(f"Generating {num_items} {worksheet_type} items and answers on: '{topic}'")
    
    # New prompt that asks for both questions and a clearly separated answer key
    prompt = (
        f"Generate a student worksheet with {num_items} {worksheet_type} questions on the topic of '{topic}'. "
        "The questions should be clear and suitable for a student. Do not provide the answers immediately after each question. "
        "After the questions, add a separator line exactly like this: '--- ANSWERS ---'. "
        "Then, after the separator, provide a numbered list of the corresponding answers. "
        "The final output must contain both the questions and the answers, separated by the '--- ANSWERS ---' line."
    )

    try:
        model = genai.GenerativeModel('gemini-1.5-flash')
        chat = model.start_chat(history=history)
        response = chat.send_message(prompt)
        return {"result": response.text, "source": "generated_worksheet_text"}
    except Exception as e:
        print(f"Error generating worksheet content: {e}")
        return {"result": "Sorry, I couldn't generate the worksheet content.", "source": "error"}
    
def generate_ppt_content(topic):
    prompt = f"Create the content for a 10-slide presentation on the topic of '{topic}'. The first slide should be a title slide. The next 9 slides should be content slides, each with a title and 3-4 bullet points. Format the output clearly, using 'SLIDE:' to mark the beginning of each slide."
    try:
        model = genai.GenerativeModel('gemini-1.5-flash')
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        print(f"Error generating PPT content: {e}")
        return None
    
def create_ppt_file(title, content):
    try:
        prs = Presentation()
        # More robustly split the content into slides
        slides = re.split(r'\n*SLIDE:?\s*\d*\s*:?\n*', content)
        slides = [s.strip() for s in slides if s.strip()] # Clean up empty entries

        if not slides:
            print("Could not parse any slides from the generated content.")
            return None

        # Handle the first slide as the title slide
        title_slide_content = slides[0].split('\n')
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_slide_content[0]
        if len(title_slide_content) > 1:
            slide.placeholders[1].text = "\n".join(title_slide_content[1:])
        
        # Handle the rest as content slides
        for slide_content in slides[1:]:
            lines = [line.strip() for line in slide_content.strip().split('\n') if line.strip()]
            if not lines: continue

            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = lines[0] # Assume first line is the title
            
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
            
            for point in lines[1:]:
                # Remove leading bullet points like '*' or '-'
                cleaned_point = re.sub(r'^\s*[\*\-]\s*', '', point)
                p = tf.add_paragraph()
                p.text = cleaned_point
                p.level = 0

        fp = BytesIO()
        prs.save(fp)
        fp.seek(0)
        return fp.read()
    except Exception as e:
        print(f"Error creating PPT file: {e}")
        traceback.print_exc()
        return None



def send_whatsapp_ppt(to, file_bytes, filename):
    media_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/media"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    files = {'file': (filename, file_bytes, 'application/vnd.openxmlformats-officedocument.presentationml.presentation'), 'messaging_product': (None, 'whatsapp')}
    try:
        media_res = requests.post(media_url, headers=headers, files=files)
        media_res.raise_for_status()
        media_id = media_res.json().get("id")
        if not media_id: return "failure"
        
        msg_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
        msg_headers = {"Authorization": f"Bearer {ACCESS_TOKEN}", "Content-Type": "application/json"}
        msg_payload = {"messaging_product": "whatsapp", "to": to, "type": "document", "document": {"id": media_id, "filename": filename}}
        requests.post(msg_url, headers=msg_headers, json=msg_payload).raise_for_status()
        return "success"
    except Exception as e:
        print(f"Failed to send PPT: {e}")
        return "failure"

def query_os_subject(question):
    try:
        embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        db = FAISS.load_local("vector_index/os_docs", embeddings, allow_dangerous_deserialization=True)
        retriever = db.as_retriever(search_kwargs={"k": 1})
        docs = retriever.invoke(question)
        if not docs or all(len(doc.page_content.strip()) == 0 for doc in docs):
            return {"result": "No relevant content found in OS database.", "source": "fallback_needed"}
        context = "\n\n".join([doc.page_content for doc in docs])
        sources = [os.path.basename(doc.metadata['source']) for doc in docs if 'source' in doc.metadata]
        source_info = "\n\n(Source: " + ", ".join(list(set(sources))) + ")" if sources else ""
        prompt = f"Based on the following context, answer the question accurately and concisely. If the context does not contain enough information, state that directly:\n\nContext:\n{context}\n\nQuestion: {question}\n\nAnswer:"
        model = genai.GenerativeModel('gemini-1.5-flash')
        response = model.generate_content(prompt)
        final_answer = response.text + source_info
        return {"result": final_answer, "source": "db"}
    except Exception as e:
        print("Error during OS subject query:", e)
        traceback.print_exc()
        return {"result": "Error during database lookup.", "source": "fallback_needed"}
    
def build_os_vector_index():
    print("Loading documents for OS vector index...")
    loaders = [
        UnstructuredFileLoader("data/OS_Syllabus.pdf"),
        UnstructuredFileLoader("data/OS_Unit_I_Overview.pptx"),
        UnstructuredFileLoader("data/OS_Unit_II_Process_Management.pptx"),
        UnstructuredFileLoader("data/OS_Unit_II_Process_Scheduling.pptx"),
        UnstructuredFileLoader("data/OS_Unit_III_Deadlocks.pptx"),
        UnstructuredFileLoader("data/OS_Unit_IV_Memory_management.pptx"),
        UnstructuredFileLoader("data/OS_Unit_IV_Virtual_Memory.pptx"),
        UnstructuredFileLoader("data/OS_Unit_V_File_System_Implementation.pptx"),
        UnstructuredFileLoader("data/OS_Unit_V_Storage_management.pptx"),
    ]
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    documents = []
    for loader in loaders:
        try:
            raw_docs = loader.load()
            split_docs = text_splitter.split_documents(raw_docs)
            documents.extend(split_docs)
            print(f"Loaded and split {len(raw_docs)} documents from {loader.file_path}")
        except Exception as e:
            print(f"Error loading {loader.file_path}: {e}")
    if not documents:
        print("No documents loaded, FAISS index will be empty.")
        return
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    db = FAISS.from_documents(documents, embeddings)
    db.save_local("vector_index/os_docs")
    print("OS Vector Index built successfully.")

def query_gemini(question, history):
    try:
        model = genai.GenerativeModel('gemini-1.5-flash')
        chat = model.start_chat(history=history)
        response = chat.send_message(question)
        return {"result": response.text}
    except Exception as e:
        print("Gemini general query error:", e)
        return {"result": "Sorry, I encountered an issue while processing your general question."}


def create_pdf_locally(title, content):
    """
    Creates a PDF file from a title and content, saving it to a temporary file
    and returning the raw bytes.
    """
    pdf_bytes = None
    temp_pdf_path = None
    try:
        # Create a temporary file to save the PDF to disk first
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_file:
            temp_pdf_path = temp_file.name

        # Initialize the PDF document
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)
        
        # Add the title
        pdf.set_font("Helvetica", 'B', 16)
        pdf.multi_cell(0, 10, title, align='C')
        pdf.ln(10)
        
        # Add the main content
        pdf.set_font("Helvetica", size=11)
        # Clean the content to handle special characters that FPDF might not support
        cleaned_content = content.encode('latin-1', 'replace').decode('latin-1')
        pdf.multi_cell(0, 7, cleaned_content)
        
        # Save the PDF to the temporary file path
        pdf.output(temp_pdf_path)

        # Read the bytes from the saved file to ensure it's not empty
        with open(temp_pdf_path, 'rb') as f:
            pdf_bytes = f.read()

        return pdf_bytes
        
    except Exception as e:
        print(f"Error creating PDF: {e}")
        return None
    finally:
        # Clean up the temporary file from disk after we're done with it
        if temp_pdf_path and os.path.exists(temp_pdf_path):
            os.remove(temp_pdf_path)

#
# ---------------------------------
# --- CORE MESSAGE PROCESSOR


def handle_final_classroom_post(sender, title):
    """Helper function to handle the final steps of posting to Classroom."""
    try:
        post_choice = user_temp_data[sender]['post_choice']
        questions = user_temp_data[sender]['questions_text']
        answers = user_temp_data[sender]['answers_text']
        
        pdf_content = questions
        pdf_filename = f"{title.replace(' ', '_')}_worksheet.pdf"
        if post_choice.lower() == "post with answers":
            pdf_content += f"\n\n--- ANSWERS ---\n{answers}"
            pdf_filename = f"{title.replace(' ', '_')}_with_answers.pdf"

        pdf_bytes = create_pdf_locally(title, pdf_content)
        
        if pdf_bytes:
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_file:
                temp_file.write(pdf_bytes)
                temp_file_path = temp_file.name
            
            drive_file_id = upload_file_to_drive(temp_file_path)
            os.remove(temp_file_path)

            if drive_file_id:
                result = post_announcement(CLASSROOM_COURSE_ID, title, drive_file_id)
                send_whatsapp_message(sender, result['result'])
            else:
                send_whatsapp_message(sender, "I created the PDF, but failed to upload it to Google Drive.")
        else:
            send_whatsapp_message(sender, "Sorry, I failed to create the final PDF for upload.")

    except KeyError:
        send_whatsapp_message(sender, "An error occurred. I've lost the worksheet details. Please start over.")
    finally:
        if sender in user_states: del user_states[sender]
        if sender in user_temp_data: del user_temp_data[sender]

def query_gemini_vision(prompt, image_bytes):
    """Sends a prompt and an image to the Gemini Vision model."""
    print("Querying Gemini Vision...")
    try:
        image_parts = [{"mime_type": "image/jpeg", "data": base64.b64encode(image_bytes).decode('utf-8')}]
        prompt_parts = [prompt] + image_parts
        
        model = genai.GenerativeModel('gemini-1.5-flash')
        response = model.generate_content(prompt_parts)
        
        if response.parts:
            return {"result": response.text, "source": "gemini_vision"}
        else:
            return {"result": "I'm sorry, I couldn't understand the image.", "source": "error"}
            
    except Exception as e:
        print(f"Error querying Gemini Vision: {e}")
        return {"result": "Sorry, there was an error processing the image.", "source": "error"}
    
def send_whatsapp_audio(to, audio_bytes, filename):
    media_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/media"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    files = {'file': (filename, audio_bytes, 'audio/mpeg'), 'messaging_product': (None, 'whatsapp')}
    try:
        media_res = requests.post(media_url, headers=headers, files=files)
        media_res.raise_for_status()
        media_id = media_res.json().get("id")
        if not media_id: return "failure"

        msg_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
        msg_payload = {"messaging_product": "whatsapp", "to": to, "type": "audio", "audio": {"id": media_id}}
        requests.post(msg_url, headers=headers, json=msg_payload).raise_for_status()
        return "success"
    except Exception as e:
        print(f"Failed to send audio: {e}")
        return "failure"
    
def generate_voiceover(text):
    print("Generating voiceover...")
    try:
        tts = gTTS(text=text, lang='en')
        fp = BytesIO()
        tts.write_to_fp(fp)
        fp.seek(0)
        return fp.read()
    except Exception as e:
        print(f"Error generating voiceover: {e}")
        return None
    
def process_message(sender, text):
    append_to_memory(sender, "user", text)
    current_state = user_states.get(sender)

    if text.lower().strip() in ['cancel', 'stop', 'menu', 'start', 'exit']:
        if current_state:
            del user_states[sender]; del user_temp_data[sender]
            send_whatsapp_message(sender, "Okay, I've canceled the current operation.")
        menu_text = "Hello! I'm your OS teaching assistant. Please choose an option:"
        options = ["Ask an OS Question", "Create a Worksheet", "Create a Video", "Create a PPT"]
        send_menu_message(sender, menu_text, options)
        append_to_memory(sender, "assistant", "Displayed main menu.")
        return

    # --- State-Based Conversation Flow ---
    if current_state == "awaiting_ppt_topic":
        try:
            topic = text.strip()
            send_whatsapp_message(sender, f"Okay, generating a 10-slide presentation on '{topic}'. This may take a moment...")
            
            ppt_content = generate_ppt_content(topic)
            if ppt_content:
                ppt_bytes = create_ppt_file(topic, ppt_content)
                if ppt_bytes:
                    status = send_whatsapp_ppt(sender, ppt_bytes, f"{topic.replace(' ', '_')}.pptx")
                    if status != "success":
                        send_whatsapp_message(sender, "I created the PPT, but there was an error sending it.")
                else:
                    send_whatsapp_message(sender, "Sorry, I generated the content but failed to create the PPT file.")
            else:
                send_whatsapp_message(sender, "Sorry, I couldn't generate content for that topic.")
        except Exception as e:
             print(f"Error in PPT generation flow: {e}")
             send_whatsapp_message(sender, "An unexpected error occurred.")
        finally:
            user_states.pop(sender, None)
            user_temp_data.pop(sender, None)
        return

    # --- State-Based Conversation Flow ---
    if current_state == "awaiting_worksheet_topic":
        if text.lower() == "other topics":
            user_states[sender] = "awaiting_custom_topic"
            send_whatsapp_message(sender, "Please type the custom topic for your worksheet.")
            return
        user_temp_data[sender]['topic'] = text
        user_states[sender] = "awaiting_worksheet_format"
        send_interactive_message(sender, f"Great! Topic is '{text}'.\nWhat format would you like?", ["PDF Worksheet", "Google Form Quiz"])
        return

    if current_state == "awaiting_custom_topic":
        user_temp_data[sender]['topic'] = text
        user_states[sender] = "awaiting_worksheet_format"
        send_interactive_message(sender, f"Great! Topic is '{text}'.\nWhat format would you like?", ["PDF Worksheet", "Google Form Quiz"])
        return
        
    if current_state == "awaiting_worksheet_format":
        user_temp_data[sender]['format'] = text
        user_states[sender] = "awaiting_worksheet_quantity"
        send_interactive_message(sender, f"Perfect, a {text}.\nHow many questions?", ["5", "10", "15"])
        return

    if current_state == "awaiting_worksheet_quantity":
        try:
            user_temp_data[sender]['quantity'] = int(text.strip())
            if user_temp_data[sender]['format'] == "Google Form Quiz":
                user_temp_data[sender]['type'] = "mcq"
                topic = user_temp_data[sender]['topic']
                quantity = user_temp_data[sender]['quantity']
                send_whatsapp_message(sender, f"Okay! Generating a {quantity}-question Google Form quiz on '{topic}'. Please wait...")
                
                worksheet_content_result = generate_worksheet_content_text(topic, quantity, "mcq", user_memory[sender])
                if worksheet_content_result.get("source") == "generated_worksheet_text":
                    form_result = create_google_form_mcq(f"Quiz: {topic}", worksheet_content_result["result"])
                    send_whatsapp_message(sender, form_result["result"])
                else:
                    send_whatsapp_message(sender, "Sorry, I couldn't generate the quiz content.")
                
                del user_states[sender]; del user_temp_data[sender]
                return
            else:
                user_states[sender] = "awaiting_worksheet_type"
                send_interactive_message(sender, f"Perfect, {text} questions.\nNow, what type of questions?", ["MCQ", "Short Answer", "Numerical"])
        except ValueError:
            send_whatsapp_message(sender, "Please select a valid number from the buttons.")
        return

    if current_state == "awaiting_worksheet_type":
        try:
            topic = user_temp_data[sender]['topic']
            quantity = user_temp_data[sender]['quantity']
            q_type = text.lower().replace(" answer", "")
            valid_types = ["mcq", "short", "numerical"]
            if q_type not in valid_types:
                send_whatsapp_message(sender, "That's not a valid type."); return

            send_whatsapp_message(sender, f"Okay! Generating {quantity} {q_type} questions on '{topic}'. Please wait...")
            
            worksheet_content_result = generate_worksheet_content_text(topic, quantity, q_type, user_memory[sender])
            
            if worksheet_content_result and worksheet_content_result.get("source") == "generated_worksheet_text":
                full_content = worksheet_content_result["result"]
                questions_text, answers_text = (full_content.split("--- ANSWERS ---", 1) + ["No answer key generated."])[:2]
                user_temp_data[sender]['questions_text'] = questions_text.strip()
                user_temp_data[sender]['answers_text'] = answers_text.strip()
                
                if worksheet_pdf_bytes := create_pdf_locally(f"Worksheet: {topic.title()}", questions_text.strip()):
                    send_whatsapp_document(sender, worksheet_pdf_bytes, f"{topic.replace(' ', '_')}_worksheet.pdf")
                
                if answer_key_pdf_bytes := create_pdf_locally(f"Answer Key: {topic.title()}", answers_text.strip()):
                    send_whatsapp_document(sender, answer_key_pdf_bytes, f"{topic.replace(' ', '_')}_answers.pdf")

                user_states[sender] = "awaiting_classroom_post_choice"
                send_interactive_message(sender, "I've sent the PDFs. Post to Google Classroom?", ["Post Questions Only", "Post with Answers", "Don't Post"])
            else:
                send_whatsapp_message(sender, "Sorry, I couldn't generate the worksheet content.")
                del user_states[sender]; del user_temp_data[sender]
        except (ValueError, KeyError):
            send_whatsapp_message(sender, "An error occurred. Please start over.")
            if sender in user_states: del user_states[sender]
            if sender in user_temp_data: del user_temp_data[sender]
        return
        
    if current_state == "awaiting_classroom_post_choice":
        if text.lower() == "don't post":
            send_whatsapp_message(sender, "Okay, I won't post it to Classroom. Let me know if you need anything else!")
            del user_states[sender]; del user_temp_data[sender]
            return
        
        user_temp_data[sender]['post_choice'] = text
        user_states[sender] = "awaiting_classroom_title_choice"
        topic = user_temp_data[sender].get('topic', 'Worksheet')
        default_title = f"{topic} Tutorial"
        if len(default_title) > 20:
            default_title = f"{topic} Notes"
        send_interactive_message(sender, "Great! What should be the title?", [default_title, "Type Manually"])
        return

    if current_state == "awaiting_classroom_title_choice":
        if text.lower() == "type manually":
            user_states[sender] = "awaiting_custom_classroom_title"
            send_whatsapp_message(sender, "Please type the title for the announcement.")
            return
        else:
            handle_final_classroom_post(sender, text)
            return

    if current_state == "awaiting_custom_classroom_title":
        handle_final_classroom_post(sender, text.strip())
        return
        
    if current_state == "awaiting_video_topic":
        topic = text.strip()
        send_whatsapp_message(sender, f"Great! I will start generating a video about '{topic}'. This might take a few moments.")
        # In a real implementation, you would trigger your video generation functions here.
        del user_states[sender]; del user_temp_data[sender]
        return

    # --- Initial Triggers for New Conversations ---
    if text.lower() == "create a worksheet":
        user_states[sender] = "awaiting_worksheet_topic"
        send_interactive_message(sender, "Let's create a worksheet! Please choose a topic:", ["CPU Scheduling", "Deadlocks", "Other Topics"])
        return
    
    if text.lower() == "create a ppt":
        user_states[sender] = "awaiting_ppt_topic"
        send_whatsapp_message(sender, "Excellent! What topic would you like the presentation to be about?")
        return
        
    if text.lower() == "create a video":
        send_whatsapp_message(sender, "This feature is coming soon!")
        return
    
    if text.lower() == "ask os question":
        send_whatsapp_message(sender, "Of course! What is your question about Operating Systems?")
        return
        
    # --- Fallback to a General Query ---
    print(f"Handling as a general query: '{text}'")
    response = query_os_subject(text) # This should be your smart RAG/Gemini function
    send_whatsapp_message(sender, response["result"])
    append_to_memory(sender, "assistant", response["result"])


def parse_deadline(text):
    text = text.lower()
    today = datetime.now()
    if "tomorrow" in text:
        return today + timedelta(days=1)
    if "next week" in text or "7 days" in text:
        return today + timedelta(days=7)
   
    return today + timedelta(days=2)


def send_interactive_message(to, text, buttons):
    """Sends an interactive message with up to 3 reply buttons."""
    url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}", "Content-Type": "application/json"}
    
    button_payload = []
    for i, button_title in enumerate(buttons):
        button_payload.append({
            "type": "reply",
            "reply": {
                "id": f"button_{i+1}", # A unique ID for the button
                "title": button_title
            }
        })

    payload = {
        "messaging_product": "whatsapp",
        "to": to,
        "type": "interactive",
        "interactive": {
            "type": "button",
            "body": {
                "text": text
            },
            "action": {
                "buttons": button_payload
            }
        }
    }
    
    try:
        res = requests.post(url, headers=headers, json=payload)
        res.raise_for_status()
        print(f"Interactive message sent to {to}.")
    except Exception as e:
        print(f"Failed to send interactive message: {e}")

    



#
# ---------------------------------
# --- FLASK WEBHOOK ENDPOINTS (Corrected Version)
# ---------------------------------
@app.route("/webhook", methods=["GET", "POST"])
def webhook():
    if request.method == "GET":
        if request.args.get("hub.verify_token") == VERIFY_TOKEN:
            return request.args.get("hub.challenge")
        return "Verification token mismatch", 403
    
    data = request.get_json()
    try:
        if data and "entry" in data:
            for entry in data.get("entry", []):
                for change in entry.get("changes", []):
                    value = change.get("value", {})
                    if "messages" in value:
                        for msg in value.get("messages", []):
                            sender = msg["from"]
                            text = ""
                            msg_type = msg.get("type")

                            if msg_type == "text":
                                text = msg["text"]["body"]
                            elif msg_type == "interactive":
                                interactive_type = msg["interactive"].get("type")
                                if interactive_type == "button_reply":
                                    text = msg["interactive"]["button_reply"]["title"]
                                elif interactive_type == "list_reply":
                                    text = msg["interactive"]["list_reply"]["title"]
                            elif msg_type == "button":
                                text = msg["button"]["text"]
                            elif msg_type == "audio":
                                media_id = msg['audio']['id']
                                thread = threading.Thread(target=handle_audio_task, args=(sender, media_id))
                                thread.start()
                                return "ok", 200
                            elif msg_type == "image":
                                prompt = msg.get("image", {}).get("caption", "Explain this image.")
                                media_id = msg['image']['id']
                                thread = threading.Thread(target=handle_image_task, args=(sender, media_id, prompt))
                                thread.start()
                                return "ok", 200
                            
                            if text:
                                process_message(sender, text)
    except Exception as e:
        print(f"Webhook processing error: {e}")
        traceback.print_exc()
    return "ok", 200

@app.route("/status", methods=["GET"])
def status():
    return {"status": "running"}, 200

if __name__ == "__main__":
    os.makedirs("data", exist_ok=True)
    os.makedirs("vector_index", exist_ok=True)
    init_google_apis()
    
    # --- Send Startup Template Message ---
    print("Sending startup template to users...")
    for number in student_phone_numbers:
        send_start_template(number)
    print("Finished sending templates.")

    # Create embeddings for the intent classification examples
    print("Embedding example questions for classification...")
    os_example_embeddings = np.array(embeddings_model_for_classification.embed_documents(OS_EXAMPLES))
    worksheet_generation_example_embeddings = np.array(embeddings_model_for_classification.embed_documents(WORKSHEET_GENERATION_EXAMPLES))
    # ... (add your other embedding lines here) ...
    print("Embeddings created successfully.")

    # Check if the local vector database exists, if not, build it
    if not os.path.exists("vector_index/os_docs"):
        print("Building OS vector index for the first time...")
        build_os_vector_index()
    else:
        print("OS Vector Index already exists. Skipping build.")

    # ✅ This is the crucial part that starts the server
    print("Starting Flask app...")
    app.run(port=5000, debug=False)